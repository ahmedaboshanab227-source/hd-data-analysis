import streamlit as st
import pandas as pd
import plotly.express as px
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import io

# 1. Page Configuration & Theme Styling
st.set_page_config(page_title="EXTRA HD DATA ANALYZER", layout="wide")

# Advanced Custom CSS for 3D Neumorphic / Skeuomorphic Look
st.markdown("""
    <style>
    /* Background Gradient */
    .stApp {
        background: linear-gradient(135deg, #f0f4f8 0%, #e2ebf0 100%);
    }
    
    /* 3D Soft Embossed Headers */
    h1 {
        color: #003366 !important;
        font-family: 'Segoe UI', Tahoma, Geneva, Verdana, sans-serif;
        font-weight: 800;
        text-shadow: 2px 2px 4px rgba(0, 0, 0, 0.1), -1px -1px 0px rgba(255, 255, 255, 0.9);
    }
    h2, h3 { 
        color: #003366 !important; 
        font-family: 'Segoe UI', sans-serif;
        font-weight: 600;
    }
    
    /* 3D Glass & Neumorphic Cards for Data & File Uploader */
    div[data-testid="stFileUploader"], .stDataFrame, div[data-testid="stMetricBlock"] {
        background: #f0f4f8 !important;
        border-radius: 16px !important;
        box-shadow: 9px 9px 16px rgba(163, 177, 198, 0.6), -9px -9px 16px rgba(255, 255, 255, 0.8) !important;
        border: 1px solid rgba(255, 255, 255, 0.5) !important;
        padding: 20px !important;
    }

    /* Tabs Styling 3D buttons effect */
    button[data-baseweb="tab"] {
        background: #f0f4f8 !important;
        border-radius: 10px 10px 0px 0px !important;
        box-shadow: 3px -3px 6px rgba(163, 177, 198, 0.4), -3px -3px 6px rgba(255, 255, 255, 0.7) !important;
        border: none !important;
        margin-right: 5px !important;
        color: #003366 !important;
        font-weight: bold !important;
    }
    button[data-baseweb="tab"][aria-selected="true"] {
        background: #003366 !important;
        color: #ffffff !important;
        box-shadow: inset 2px 2px 5px rgba(0,0,0,0.3) !important;
    }

    /* 3D Hover Effects on Buttons */
    div.stButton > button {
        background: linear-gradient(145deg, #003366, #002244) !important;
        color: white !important;
        border-radius: 12px !important;
        border: none !important;
        box-shadow: 5px 5px 10px rgba(163, 177, 198, 0.6), -5px -5px 10px rgba(255, 255, 255, 0.8) !important;
        transition: all 0.2s ease-in-out !important;
        font-weight: bold !important;
        font-size: 1.1em !important;
        padding: 12px 24px !important;
    }
    div.stButton > button:hover {
        transform: translateY(-2px) !important;
        box-shadow: 7px 7px 14px rgba(163, 177, 198, 0.7), -7px -7px 14px rgba(255, 255, 255, 0.9) !important;
        background: linear-gradient(145deg, #004488, #003366) !important;
    }
    div.stButton > button:active {
        transform: translateY(1px) !important;
        box-shadow: inset 3px 3px 6px rgba(0,0,0,0.4) !important;
    }
    </style>
    """, unsafe_allow_html=True)

# 3D Stylized Header Component (No image file needed)
st.markdown("""
    <div style="text-align: center; margin-bottom: 40px; padding: 20px;">
        <h1 style="font-size: 3.8em; letter-spacing: 2px; margin-bottom: 0px; display: inline-block;">
            <span style="color: #003366; text-shadow: 3px 3px 0px #e2ebf0, 6px 6px 0px rgba(0,0,0,0.15);">E<span style="color: #FFCC00; font-size: 1.2em; font-style: italic; text-shadow: 3px 3px 0px #003366;">X</span>TRA</span> 
            <span style="font-size: 0.7em; color: #003366; font-weight: 300;">HD DATA ANALYZER</span>
        </h1>
        <p style="color: #556677; font-size: 1.2em; margin-top: 15px; font-weight: 500;">
            🚀 Advanced Logistics Intelligence & Automated Performance Insights
        </p>
    </div>
    """, unsafe_allow_html=True)

# 2. File Upload Zone
uploaded_file = st.file_uploader("📥 Drop your Excel Logistics File here (.xlsx)", type=["xlsx"])

# Helper function to generate standardized 3D-looking charts
def create_bar_chart(df, x_col, y_col, title):
    fig = px.bar(df, x=x_col, y=y_col, title=title, text_auto=True,
                 color_discrete_sequence=['#003366']) # Deep Blue Bars
    fig.update_layout(
        plot_bgcolor='rgba(0,0,0,0)',
        paper_bgcolor='rgba(0,0,0,0)',
        xaxis_title=x_col,
        yaxis_title="Unique Orders Count",
        font=dict(color="#003366", family="Segoe UI"),
        title_font=dict(size=18, color="#003366")
    )
    fig.update_traces(
        marker_line_color='#FFCC00', 
        marker_line_width=1.5,
        textposition='outside', # Puts numbers clearly on top of bars
        shadow=dict(color="rgba(0,0,0,0.3)", blur=5, x=2, y=2) # Simulated 3D bar shadow
    )
    return fig

# 3. Processing Core
if uploaded_file is not None:
    try:
        xls = pd.ExcelFile(uploaded_file)
        sheets_in_file = xls.sheet_names
        
        required_sheets = ['HD', 'Confirmation', 'Return']
        missing_sheets = [s for s in required_sheets if s not in sheets_in_file]
        
        if missing_sheets:
            st.error(f"⚠️ Error: Missing required sheets: {', '.join(missing_sheets)}")
        else:
            ppt_data = {}

            # Create Modern Navigation Tabs
            tab1, tab2, tab3 = st.tabs(["📊 HD Analysis Dashboard", "✅ Confirmation Summary", "🔄 Return Trends"])

            # ----------------------------------------
            # TAB 1: HD SHEET ANALYSIS
            # ----------------------------------------
            with tab1:
                st.markdown("<h3 style='margin-top:15px;'>🔹 High-Definition Logistics Overview</h3>", unsafe_allow_html=True)
                df_hd = pd.read_excel(xls, sheet_name='HD')
                df_hd.columns = df_hd.columns.str.strip()
                
                if 'BOOK ID' in df_hd.columns:
                    if 'Actual Delivery Date' in df_hd.columns:
                        df_hd['Actual Delivery Date'] = pd.to_datetime(df_hd['Actual Delivery Date'], errors='coerce')
                        df_hd['Month'] = df_hd['Actual Delivery Date'].dt.strftime('%Y-%m ( %B )')

                    columns_to_analyze = [
                        'Group Name', 'Region Name', 'Store Code', 'Area Name', 
                        'Technician', 'Driver', 'Truck No', 'Month'
                    ]
                    
                    ppt_data['HD'] = {}
                    
                    for col in columns_to_analyze:
                        if col in df_hd.columns or (col == 'Month' and 'Month' in df_hd.columns):
                            display_name = "Actual Delivery Date (By Month)" if col == 'Month' else col
                            
                            df_temp = df_hd.dropna(subset=[col])
                            summary = df_temp.groupby(col)['BOOK ID'].nunique().reset_index()
                            summary.columns = [display_name, 'Unique Orders']
                            summary = summary.sort_values(by='Unique Orders', ascending=False)
                            
                            ppt_data['HD'][display_name] = summary
                            
                            chart_col, table_col = st.columns([2, 1])
                            with chart_col:
                                st.plotly_chart(create_bar_chart(summary, display_name, 'Unique Orders', f"Distribution by {display_name}"), use_container_width=True)
                            with table_col:
                                st.markdown(f"<p style='color:#003366; font-weight:bold; margin-bottom:5px;'>📋 {display_name} Metrics Data:</p>", unsafe_allow_html=True)
                                st.dataframe(summary, hide_index=True, use_container_width=True)
                                
                            st.divider()
                else:
                    st.error("Column 'BOOK ID' not found in HD sheet.")

            # ----------------------------------------
            # TAB 2: CONFIRMATION SHEET ANALYSIS
            # ----------------------------------------
            with tab2:
                st.markdown("<h3 style='margin-top:15px;'>🔹 Execution & Confirmation Intelligence</h3>", unsafe_allow_html=True)
                df_conf = pd.read_excel(xls, sheet_name='Confirmation')
                df_conf.columns = df_conf.columns.str.strip()
                
                if 'BOOK ID' in df_conf.columns:
                    if 'Date' in df_conf.columns:
                        df_conf['Date'] = pd.to_datetime(df_conf['Date'], errors='coerce')
                        df_conf['Month'] = df_conf['Date'].dt.strftime('%Y-%m ( %B )')
                        
                    conf_columns = ['Technician', 'Driver', 'Truck No', 'Month']
                    ppt_data['Confirmation'] = {}
                    
                    for col in conf_columns:
                        if col in df_conf.columns or (col == 'Month' and 'Month' in df_conf.columns):
                            display_name = "Date (By Month)" if col == 'Month' else col
                            
                            df_temp = df_conf.dropna(subset=[col])
                            summary = df_temp.groupby(col)['BOOK ID'].nunique().reset_index()
                            summary.columns = [display_name, 'Unique Orders']
                            summary = summary.sort_values(by='Unique Orders', ascending=False)
                            
                            ppt_data['Confirmation'][display_name] = summary
                            
                            chart_col, table_col = st.columns([2, 1])
                            with chart_col:
                                st.plotly_chart(create_bar_chart(summary, display_name, 'Unique Orders', f"Performance by {display_name}"), use_container_width=True)
                            with table_col:
                                st.markdown(f"<p style='color:#003366; font-weight:bold; margin-bottom:5px;'>📋 {display_name} Metrics Data:</p>", unsafe_allow_html=True)
                                st.dataframe(summary, hide_index=True, use_container_width=True)
                                
                            st.divider()
                else:
                    st.error("Column 'BOOK ID' not found in Confirmation sheet.")

            # ----------------------------------------
            # TAB 3: RETURN SHEET ANALYSIS
            # ----------------------------------------
            with tab3:
                st.markdown("<h3 style='margin-top:15px;'>🔹 Operational Reverse Logistics (Returns)</h3>", unsafe_allow_html=True)
                df_ret = pd.read_excel(xls, sheet_name='Return')
                df_ret.columns = df_ret.columns.str.strip()
                
                date_col = 'Date' if 'Date' in df_ret.columns else ([c for c in df_ret.columns if 'date' in c.lower()] + [None])[0]
                
                if 'BOOK ID' in df_ret.columns and date_col:
                    df_ret['Parsed Date'] = pd.to_datetime(df_ret[date_col], errors='coerce')
                    df_ret['Month'] = df_ret['Parsed Date'].dt.strftime('%Y-%m ( %B )')
                    
                    df_temp = df_ret.dropna(subset=['Month'])
                    summary_ret = df_temp.groupby('Month')['BOOK ID'].nunique().reset_index()
                    summary_ret.columns = ['Month', 'Unique Returns']
                    summary_ret = summary_ret.sort_values(by='Month')
                    
                    ppt_data['Return'] = summary_ret
                    
                    chart_col, table_col = st.columns([2, 1])
                    with chart_col:
                        st.plotly_chart(create_bar_chart(summary_ret, 'Month', 'Unique Returns', "Returned Orders Trend Analysis"), use_container_width=True)
                    with table_col:
                        st.markdown("<p style='color:#003366; font-weight:bold; margin-bottom:5px;'>📋 Returns Trend Data:</p>", unsafe_allow_html=True)
                        st.dataframe(summary_ret, hide_index=True, use_container_width=True)
                else:
                    st.error("Make sure 'Return' sheet has both 'BOOK ID' and a valid 'Date' column.")

            # ----------------------------------------
            # POWERPOINT EXPORT LOGIC
            # ----------------------------------------
            st.markdown("<br><h2 style='text-align: center;'>💼 Executive Report Deployment</h2>", unsafe_allow_html=True)
            
            if st.button("⚡ Compiling Executive PowerPoint Report 📊", use_container_width=True):
                prs = Presentation()
                DARK_BLUE = RGBColor(0, 51, 102)
                YELLOW = RGBColor(255, 204, 0)
                
                # Title Slide
                slide = prs.slides.add_slide(prs.slide_layouts[6])
                tx_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(2))
                tf = tx_box.text_frame
                p = tf.add_paragraph()
                p.text = "EXTRA HD DATA ANALYZER"
                p.font.size = Pt(40)
                p.font.bold = True
                p.font.color.rgb = DARK_BLUE
                
                p2 = tf.add_paragraph()
                p2.text = "Automated Logistics & Performance Summary Report"
                p2.font.size = Pt(18)
                p2.font.color.rgb = YELLOW
                
                for sheet_name, categories in ppt_data.items():
                    if sheet_name in ['HD', 'Confirmation']:
                        for cat_name, df_summary in categories.items():
                            slide = prs.slides.add_slide(prs.slide_layouts[6])
                            tx = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
                            tx.text_frame.text = f"{sheet_name} Sheet - {cat_name} Analysis"
                            tx.text_frame.paragraphs[0].font.size = Pt(24)
                            tx.text_frame.paragraphs[0].font.bold = True
                            tx.text_frame.paragraphs[0].font.color.rgb = DARK_BLUE
                            
                            df_top = df_summary.head(10)
                            rows = len(df_top) + 1
                            table_shape = slide.shapes.add_table(rows, 2, Inches(1), Inches(1.5), Inches(8), Inches(4))
                            table = table_shape.table
                            
                            table.cell(0, 0).text = str(df_top.columns[0])
                            table.cell(0, 1).text = str(df_top.columns[1])
                            table.cell(0, 0).fill.solid()
                            table.cell(0, 0).fill.fore_color.rgb = DARK_BLUE
                            table.cell(0, 1).fill.solid()
                            table.cell(0, 1).fill.fore_color.rgb = DARK_BLUE
                            
                            for r_idx, row in df_top.reset_index(drop=True).iterrows():
                                table.cell(r_idx + 1, 0).text = str(row.iloc[0])
                                table.cell(r_idx + 1, 1).text = str(row.iloc[1])
                                
                    elif sheet_name == 'Return':
                        slide = prs.slides.add_slide(prs.slide_layouts[6])
                        tx = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
                        tx.text_frame.text = "Return Sheet - Monthly Summary"
                        tx.text_frame.paragraphs[0].font.size = Pt(24)
                        tx.text_frame.paragraphs[0].font.bold = True
                        tx.text_frame.paragraphs[0].font.color.rgb = DARK_BLUE
                        
                        rows = len(categories) + 1
                        table_shape = slide.shapes.add_table(rows, 2, Inches(1), Inches(1.5), Inches(8), Inches(4))
                        table = table_shape.table
                        table.cell(0, 0).text = "Month"
                        table.cell(0, 1).text = "Unique Returns"
                        table.cell(0, 0).fill.solid()
                        table.cell(0, 0).fill.fore_color.rgb = DARK_BLUE
                        table.cell(0, 1).fill.solid()
                        table.cell(0, 1).fill.fore_color.rgb = DARK_BLUE
                        
                        for r_idx, row in categories.reset_index(drop=True).iterrows():
                            table.cell(r_idx + 1, 0).text = str(row.iloc[0])
                            table.cell(r_idx + 1, 1).text = str(row.iloc[1])

                ppt_buffer = io.BytesIO()
                prs.save(ppt_buffer)
                ppt_buffer.seek(0)
                
                st.download_button(
                    label="📥 Click here to save the Executive PowerPoint file",
                    data=ppt_buffer,
                    file_name="EXTRA_HD_Professional_Report.pptx",
                    mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
                )
                st.success("✨ Professional 3D Executive layout compiled successfully!")
                
    except Exception as e:
        st.error(f"❌ An error occurred while parsing the file: {e}")
